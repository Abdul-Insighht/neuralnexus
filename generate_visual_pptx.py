import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Add project root to path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx_helpers import (
    set_bg, add_decorative_bar, add_bottom_bar, add_card,
    add_title_text, add_subtitle, add_body_text, add_icon_circle,
    add_arrow_right, add_arrow_down, add_stat_box, add_table_slide,
    PARCHMENT, DARK_INK, WARM_BROWN, DEEP_BLUE, ACCENT_GOLD,
    ACCENT_RED, ACCENT_GREEN, ACCENT_TEAL, LIGHT_BOX, WHITE,
    CARD_BG, BORDER_COLOR
)

def build_presentation():
    prs = Presentation()
    # Use 16:9 widescreen layout standard (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # We will use slide_layouts[6] which is a completely blank layout,
    # so we can manually place shapes and control the exact visual hierarchy.
    blank_layout = prs.slide_layouts[6]
    
    # ── SLIDE 1: TITLE COVER ────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    
    # Left border band (Gold)
    gold_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), Inches(7.5))
    gold_band.fill.solid()
    gold_band.fill.fore_color.rgb = ACCENT_GOLD
    gold_band.line.fill.background()
    
    # Title & Subtitle in a single TextFrame for perfect flow
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.0), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "NeuralNexus"
    p1.font.name = 'Georgia'
    p1.font.size = Pt(56)
    p1.font.bold = True
    p1.font.color.rgb = DEEP_BLUE
    
    p2 = tf.add_paragraph()
    p2.text = "Autonomous Financial Intelligence Fabric"
    p2.font.name = 'Georgia'
    p2.font.size = Pt(24)
    p2.font.italic = True
    p2.font.color.rgb = WARM_BROWN
    p2.space_before = Pt(12)
    
    p3 = tf.add_paragraph()
    p3.text = "5 AI Agents that detect contradictions, reconcile conflicts, and audit every answer with full lineage."
    p3.font.name = 'Georgia'
    p3.font.size = Pt(14)
    p3.font.color.rgb = DARK_INK
    p3.space_before = Pt(18)
    
    # Project Context Panel
    context_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.2), Inches(11.3), Inches(2.4))
    context_shape.fill.solid()
    context_shape.fill.fore_color.rgb = CARD_BG
    context_shape.line.color.rgb = BORDER_COLOR
    context_tf = context_shape.text_frame
    context_tf.word_wrap = True
    context_tf.margin_left = Inches(0.2)
    context_tf.margin_top = Inches(0.2)
    
    p_tr = context_tf.paragraphs[0]
    p_tr.text = "📊 Track 4: Data & Intelligence  |  AI Agent Olympics 2026"
    p_tr.font.name = 'Georgia'
    p_tr.font.size = Pt(14)
    p_tr.font.bold = True
    p_tr.font.color.rgb = DEEP_BLUE
    
    p_team = context_tf.add_paragraph()
    p_team.text = "Team Members:  Hafiz Abdul Rehman  •  Meer Hamza  •  Sayyam Akram"
    p_team.font.name = 'Georgia'
    p_team.font.size = Pt(12)
    p_team.font.bold = True
    p_team.font.color.rgb = WARM_BROWN
    p_team.space_before = Pt(10)
    
    p_urls = context_tf.add_paragraph()
    p_urls.text = "Live App:  neuralnexus-omega.vercel.app\nGitHub:    github.com/Abdul-Insighht/neuralnexus"
    p_urls.font.name = 'Courier New'
    p_urls.font.size = Pt(11)
    p_urls.font.color.rgb = DARK_INK
    p_urls.space_before = Pt(12)
    
    # ── SLIDE 2: THE PROBLEM ────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    add_title_text(slide, "The Enterprise Data Quality Crisis")
    add_subtitle(slide, "Modern enterprises lose millions annually due to silent data discrepancies across systems.")
    add_bottom_bar(slide)
    
    # Large Stat Callout (Left side)
    add_stat_box(slide, Inches(0.8), Inches(1.8), "$12.9M", "Average Annual Loss\nDue to Poor Data Quality", ACCENT_RED)
    
    # Large Warning/Fact Box
    warn_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.3), Inches(2.2), Inches(3.4))
    warn_box.fill.solid()
    warn_box.fill.fore_color.rgb = LIGHT_BOX
    warn_box.line.color.rgb = ACCENT_GOLD
    warn_tf = warn_box.text_frame
    warn_tf.word_wrap = True
    warn_tf.margin_left = Inches(0.15)
    warn_tf.margin_top = Inches(0.15)
    p_warn = warn_tf.paragraphs[0]
    p_warn.text = "🚨 Critical Risk Factors"
    p_warn.font.bold = True
    p_warn.font.size = Pt(13)
    p_warn.font.color.rgb = ACCENT_RED
    p_warn.font.name = 'Georgia'
    
    p_warn_body = warn_tf.add_paragraph()
    p_warn_body.text = "When Sales DB disagrees with audited reports, decisions stall. Inaccurate reports lead to compliance failures, heavy regulatory penalties, and completely broken executive forecasts."
    p_warn_body.font.size = Pt(11)
    p_warn_body.font.color.rgb = DARK_INK
    p_warn_body.font.name = 'Georgia'
    p_warn_body.space_before = Pt(10)
    
    # Cards grid (Right side: 2x2 grid)
    add_card(slide, Inches(3.3), Inches(1.8), Inches(4.5), Inches(2.3),
             "1. Fragmented Data Silos",
             "Financial indicators live across 5 to 15 different applications (CRMs, SQL databases, offline spreadsheets, external APIs, and manual PDF reports) without coordination.",
             title_color=DEEP_BLUE)
             
    add_card(slide, Inches(8.0), Inches(1.8), Inches(4.5), Inches(2.3),
             "2. Silent Contradictions",
             "Numerical figures often conflict silently across platforms (e.g. Sales system shows $4.2M, Financials show $3.8M). No existing pipeline flags these discrepancies automatically.",
             title_color=DEEP_BLUE)
             
    add_card(slide, Inches(3.3), Inches(4.4), Inches(4.5), Inches(2.3),
             "3. 100% Reactive Workflows",
             "Financial intelligence analysts spend up to 80% of their operational bandwidth just locating data errors and trying to reconcile them manually, leaving no time for real risk mitigation.",
             title_color=DEEP_BLUE)
             
    add_card(slide, Inches(8.0), Inches(4.4), Inches(4.5), Inches(2.3),
             "4. Zero Verification & Traceability",
             "Standard corporate AI systems produce outputs without citations or traceability, causing hallucinations. There is no proof or lineage trace back to the exact source row and column.",
             title_color=DEEP_BLUE)
             
    # ── SLIDE 3: HACKATHON TRACK alignment ──────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    add_title_text(slide, "📊 Track 4: Data & Intelligence")
    add_subtitle(slide, "Purpose-built alignment: how NeuralNexus addresses each core challenge of Track 4.")
    add_bottom_bar(slide)
    
    headers = ["Track 4 Focus Area", "NeuralNexus Enterprise Implementation"]
    rows = [
        ["RAG Systems Over Multi-Source Data", "Temporal Vector Store using 768-dim Gemini Embeddings + automatic metadata filters to query structured (CSV, Excel) and unstructured (PDF, Images) federated data."],
        ["AI-Powered Pipelines & Validation", "ValidatorAgent runs continuous background scans (Z-score, CUSUM-like drift, spikes). DataQualityScorer grades every source dynamically on Completeness, Consistency, Timeliness."],
        ["Natural Language Querying Agents", "OrchestratorAgent decomposes queries, routes to specialized sub-agents, and streams unified answers to the frontend in real-time using NDJSON protocol."],
        ["Anomaly Detection & Forecasting", "ValidatorAgent identifies outliers & temporal pattern breaks, while ForecasterAgent computes trends, seasonal autocorrelations, and outputs proactive risk assessments."],
        ["Knowledge Graph Extraction", "DynamicKnowledgeGraph automatically parses documents & tables to construct a complete NetworkX entity network, dynamically injecting contradiction edges (red)."]
    ]
    
    add_table_slide(slide, headers, rows, left=Inches(0.8), top=Inches(1.8), col_widths=[Inches(3.3), Inches(8.4)])
    
    # ── SLIDE 4: THE SOLUTION ───────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    add_title_text(slide, "The Autonomous Intelligence Fabric")
    add_subtitle(slide, "A unified 5-agent ecosystem designed to act as an automated CFO assistant.")
    add_bottom_bar(slide)
    
    # 3x2 visual grid of solution pillars
    add_card(slide, Inches(0.8), Inches(1.8), Inches(3.6), Inches(2.3),
             "1. Multi-Format Federation Hub",
             "Ingests CSV, Excel, JSON, PDF (via Multimodal Vision), and Images (receipts, charts) into a single, unified database schema.",
             title_color=DEEP_BLUE)
             
    add_card(slide, Inches(4.7), Inches(1.8), Inches(3.6), Inches(2.3),
             "2. Automatic Conflict Finder",
             "Scans cross-source data dynamically to uncover numerical discrepancies, semantic contradictions, and temporal metric variations.",
             title_color=DEEP_BLUE)
             
    add_card(slide, Inches(8.6), Inches(1.8), Inches(3.9), Inches(2.3),
             "3. Trust-Weighted Reconciler",
             "Leverages multi-agent reasoning to compare data freshness, source credibility, and record count, recommending the logical 'best-truth' value.",
             title_color=DEEP_BLUE)
             
    add_card(slide, Inches(0.8), Inches(4.4), Inches(3.6), Inches(2.3),
             "4. Proactive Forecasting Hub",
             "Computes time-series trend strength, seasonal cycles, and proactive risk indices, delivering predictive warnings before the user asks.",
             title_color=DEEP_BLUE)
             
    add_card(slide, Inches(4.7), Inches(4.4), Inches(3.6), Inches(2.3),
             "5. Five-Dimension Answer Audit",
             "The CriticAgent acts as an independent compliance gate, grading every output on Accuracy, Completeness, Uncertainty, and Actionability.",
             title_color=DEEP_BLUE)
             
    add_card(slide, Inches(8.6), Inches(4.4), Inches(3.9), Inches(2.3),
             "6. Full Cell-Level Data Lineage",
             "Every word and number is fully traceable. Click on any figure in a generated answer to see the exact CSV row, column, or PDF sentence it came from.",
             title_color=DEEP_BLUE)
             
    # ── SLIDE 5: SYSTEM ARCHITECTURE (VISUAL FLOW) ───────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    add_title_text(slide, "Autonomous Multi-Agent Orchestration")
    add_subtitle(slide, "Visual data pipeline flow, showcasing sequential processing and verification gates.")
    add_bottom_bar(slide)
    
    # Draw sequential boxes with arrows
    # Input
    add_icon_circle(slide, Inches(0.8), Inches(3.0), Inches(1.6), DEEP_BLUE, "1. Ingestion\n---\nFederation\nHub")
    add_arrow_right(slide, Inches(2.5), Inches(3.65), Inches(0.4), Inches(0.3))
    
    # Brain / Orchestrator
    add_icon_circle(slide, Inches(3.0), Inches(2.6), Inches(2.4), WARM_BROWN, "2. Orchestrator\n---\nThe Agent Brain\nDecomposes\nQueries")
    
    # 3 Arrows branching from Orchestrator to specialists
    add_arrow_right(slide, Inches(5.5), Inches(1.85), Inches(0.5), Inches(0.2)) # Up-rightish
    add_arrow_right(slide, Inches(5.5), Inches(3.65), Inches(0.5), Inches(0.2)) # Right
    add_arrow_right(slide, Inches(5.5), Inches(5.45), Inches(0.5), Inches(0.2)) # Down-rightish
    
    # Specialist Box 1
    add_card(slide, Inches(6.1), Inches(1.4), Inches(2.4), Inches(1.3),
             "🔍 ValidatorAgent", "Statistical anomalies, outliers & CUSUM drift scanning.", bg_color=WHITE)
             
    # Specialist Box 2
    add_card(slide, Inches(6.1), Inches(3.1), Inches(2.4), Inches(1.3),
             "⚖️ ReconcilerAgent", "Trust-weighted mathematical resolving of conflicts.", bg_color=WHITE)
             
    # Specialist Box 3
    add_card(slide, Inches(6.1), Inches(4.8), Inches(2.4), Inches(1.3),
             "📈 ForecasterAgent", "Seasonality, rolling trends & proactive risk scores.", bg_color=WHITE)
             
    # Arrows merging to Critic
    add_arrow_right(slide, Inches(8.6), Inches(1.85), Inches(0.5), Inches(0.2))
    add_arrow_right(slide, Inches(8.6), Inches(3.65), Inches(0.5), Inches(0.2))
    add_arrow_right(slide, Inches(8.6), Inches(5.45), Inches(0.5), Inches(0.2))
    
    # Critic
    add_icon_circle(slide, Inches(9.2), Inches(2.8), Inches(2.0), ACCENT_RED, "3. CriticAgent\n---\nQA Gate\nScores QA &\nInjects Caveats")
    add_arrow_right(slide, Inches(11.3), Inches(3.65), Inches(0.4), Inches(0.3))
    
    # Output
    add_icon_circle(slide, Inches(11.8), Inches(3.0), Inches(1.3), ACCENT_GREEN, "4. Output\n---\nNDJSON\nStream")
    
    # ── SLIDE 6: THE 5 AI AGENTS ─────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    add_title_text(slide, "The Coordinated Agent Specialists")
    add_subtitle(slide, "Five distinct, prompt-engineered AI agents mapped to specific mathematical & analytical tasks.")
    add_bottom_bar(slide)
    
    # 5 vertical columns
    col_w = Inches(2.2)
    spacing = Inches(0.2)
    left_start = Inches(0.8)
    
    agents_info = [
        ("OrchestratorAgent", "THE BRAIN", DEEP_BLUE,
         "Coordinates the complete pipeline. Decomposes complex financial queries, triggers vector retrieval, dynamically handles context windows, and manages asynchronous tasks."),
        ("ValidatorAgent", "THE WATCHDOG", ACCENT_TEAL,
         "Continuously executes statistical scans on ingested databases. Uses Z-scores to isolate anomalies, CUSUM drift detection, and flags quality drop warnings (🔴 Critical, 🟠 Warning)."),
        ("ReconcilerAgent", "THE JUDGE", WARM_BROWN,
         "Analyzes and resolves data mismatches. Weighs source credibility (e.g. Audited File = 0.95 vs Spreadsheet = 0.70) to produce 'best truth' values with an uncertainty range."),
        ("ForecasterAgent", "THE ORACLE", ACCENT_GREEN,
         "Detects temporal patterns and trend indices. Computes Coefficient of Variation (volatility) and autocorrelation at multiple lag intervals to model future risks."),
        ("CriticAgent", "THE AUDITOR", ACCENT_RED,
         "Strict Quality Assurance firewall. Rates every orchestrated output (0-100) on Accuracy, Completeness, Actionability, and handling of Uncertainty. Automatically adds disclaimers.")
    ]
    
    for i, (name, role, color, desc) in enumerate(agents_info):
        left = left_start + i * (col_w + spacing)
        
        # Header block
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.8), col_w, Inches(1.0))
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        header_tf = header.text_frame
        header_tf.word_wrap = True
        p_name = header_tf.paragraphs[0]
        p_name.text = name
        p_name.font.size = Pt(13)
        p_name.font.bold = True
        p_name.font.color.rgb = WHITE
        p_name.font.name = 'Georgia'
        p_name.alignment = PP_ALIGN.CENTER
        
        p_role = header_tf.add_paragraph()
        p_role.text = role
        p_role.font.size = Pt(10)
        p_role.font.bold = True
        p_role.font.color.rgb = PARCHMENT
        p_role.font.name = 'Georgia'
        p_role.space_before = Pt(4)
        p_role.alignment = PP_ALIGN.CENTER
        
        # Body block
        body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.9), col_w, Inches(3.8))
        body.fill.solid()
        body.fill.fore_color.rgb = CARD_BG
        body.line.color.rgb = BORDER_COLOR
        body_tf = body.text_frame
        body_tf.word_wrap = True
        body_tf.margin_left = Inches(0.15)
        body_tf.margin_top = Inches(0.15)
        p_desc = body_tf.paragraphs[0]
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = DARK_INK
        p_desc.font.name = 'Georgia'
        
    # ── SLIDE 7: CORE INFRASTRUCTURE ─────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    add_title_text(slide, "Under-the-Hood Technical Modules")
    add_subtitle(slide, "Engineered framework layers that process data pipelines, index metadata, and build lineage.")
    add_bottom_bar(slide)
    
    # 3x2 grid of infrastructure
    infra_info = [
        ("Data Federation Hub", "Layer 1: Unified Ingestion", 
         "Standardizes CSV, Excel, JSON, PDFs, and images into normalized Pandas DataFrames. Features automatic schema inference and SHA-256 fingerprinting for secure data lifecycle control."),
        ("Data Quality Scorer", "Layer 2: Multi-Dimension Grades", 
         "Automatically scores every incoming dataset. Weights Completeness (35% null check), Consistency (35% duplicate check), and Timeliness (30% freshness index) to assign grades A to F."),
        ("Contradiction Detector", "Layer 2: Factual Conflict Engine", 
         "Cross-checks data in 3 modes: tabular vs tabular, extracted PDF text vs tabular aggregated metrics, and text vs text using regex entities. Classifies discrepancy severity dynamically."),
        ("Data Lineage Tracker", "Layer 2: Factual Provenance", 
         "Tracks cell-level, text-span-level, and AI inference data chains. Generates full traceability maps, letting the system trace any synthesized fact directly to its origin."),
        ("Dynamic Knowledge Graph", "Layer 3: Graph Network (PyVis)", 
         "Automatically extracts semantic relationships and entities. Projects them into a NetworkX structure. Renders red edge lines to mark contradictions in real time."),
        ("Temporal Vector Store", "Layer 3: FAISS Semantic Retrieval", 
         "Creates 768-dim vector embeddings of chunked text using Gemini. Implements exponential recency decay (decay factor = 0.95) to prioritize the latest financial filings.")
    ]
    
    for i, (title, layer, text) in enumerate(infra_info):
        row = i // 3
        col = i % 3
        
        left = Inches(0.8) + col * Inches(3.9)
        top = Inches(1.8) + row * Inches(2.6)
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = BORDER_COLOR
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = DEEP_BLUE
        p1.font.name = 'Georgia'
        
        p2 = tf.add_paragraph()
        p2.text = layer
        p2.font.size = Pt(10)
        p2.font.bold = True
        p2.font.color.rgb = WARM_BROWN
        p2.font.name = 'Georgia'
        p2.space_before = Pt(2)
        
        p3 = tf.add_paragraph()
        p3.text = text
        p3.font.size = Pt(11)
        p3.font.color.rgb = DARK_INK
        p3.font.name = 'Georgia'
        p3.space_before = Pt(8)

    # ── SLIDE 8: THE USER INTERFACE ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    add_title_text(slide, "The Enterprise Intelligence Interface")
    add_subtitle(slide, "A premium Bookish design system with custom CSS, interactive frames, and responsive dashboards.")
    add_bottom_bar(slide)
    
    # 3x2 grid of UI modules
    ui_modules = [
        ("1. Intelligence Query Chat", "Streams real-time answers (NDJSON fetch reader). Dynamically embeds interactive source-citation blocks and agent confidence bars."),
        ("2. Data Quality Dashboard", "Interactive card deck listing per-source grades (A-F), total cells analyzed, and detailed null/duplicate issues."),
        ("3. Contradiction Panel", "Displays flagged tabular-to-document discrepancies. Groups by high, medium, and low severity with interactive click resolution."),
        ("4. Proactive Analytics Hub", "Presents rolling time-series alerts, Coefficient of Variation percentages, and automatically derived trend risk rankings."),
        ("5. Knowledge Graph Tab", "Embeds the custom interactive PyVis physics layout. Color-codes entities and draws conflict edges in red for instant discovery."),
        ("6. Cell-Level Lineage Explorer", "Displays JSON/Table schemas. Enables deep search, resolving downstream fact nodes to the exact row and column index.")
    ]
    
    for i, (title, desc) in enumerate(ui_modules):
        row = i // 3
        col = i % 3
        
        left = Inches(0.8) + col * Inches(3.9)
        top = Inches(1.8) + row * Inches(2.6)
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = ACCENT_GOLD
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.15)
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = DEEP_BLUE
        p1.font.name = 'Georgia'
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK_INK
        p2.font.name = 'Georgia'
        p2.space_before = Pt(8)

    # ── SLIDE 9: GEMINI INTEGRATION ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    add_title_text(slide, "Deep Google Gemini API Integration")
    add_subtitle(slide, "How NeuralNexus harnesses Gemini models to construct an intelligent agent fabric.")
    add_bottom_bar(slide)
    
    headers_gemini = ["System Task", "Gemini Model / Technology Used", "Operational Pipeline Purpose"]
    rows_gemini = [
        ["Agent Core Reasoning", "Gemini 2.5 Flash Lite (Content Generation)", "Powering Orchestrator, Validator, Reconciler, Forecaster, and Critic reasoning pipelines with extremely low latency."],
        ["Semantic Retrieval", "Gemini Embedding-2 (models/gemini-embedding-2)", "Generating high-dimensional (768-dim) vectors of data chunks to facilitate contextual RAG searches."],
        ["PDF Data Extraction", "Gemini Multimodal Vision API", "Extracting tabular data and raw text structures from highly complex, multi-page PDF documents."],
        ["Image Analysis", "Gemini Multimodal Vision API", "Ingesting screenshots, invoices, physical receipts, and analytical charts to transform image facts into structured database schemas."],
        ["Real-Time Streaming", "Gemini API Stream (stream=True)", "Enabling token-by-token content generation. Frontend reads HTTP stream immediately via Fetch API for responsive UI."],
        ["Quotas & Stability", "Thread-Safe Rate Limiter Integration", "Custom middleware spacing API queries to eliminate 429 quota exceptions, guaranteeing uptime on standard models."]
    ]
    
    add_table_slide(slide, headers_gemini, rows_gemini, left=Inches(0.8), top=Inches(1.8), col_widths=[Inches(2.5), Inches(4.3), Inches(4.9)])

    # ── SLIDE 10: BUSINESS VALUE & MARKET SCOPE ──────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    add_title_text(slide, "Strategic Business Value & Market Scope")
    add_subtitle(slide, "Quantifying the commercial footprint and core enterprise monetization pathways.")
    add_bottom_bar(slide)
    
    # Stat boxes
    add_stat_box(slide, Inches(0.8), Inches(1.8), "$15.8B", "Total Addressable Market\nEnterprise Data Quality", DEEP_BLUE)
    add_stat_box(slide, Inches(0.8), Inches(3.3), "$4.2B", "Serviceable Market\nAI-Native Data Auditing", ACCENT_GOLD)
    
    # Large target box
    target_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.8), Inches(2.2), Inches(1.9))
    target_box.fill.solid()
    target_box.fill.fore_color.rgb = LIGHT_BOX
    target_box.line.color.rgb = BORDER_COLOR
    target_tf = target_box.text_frame
    target_tf.word_wrap = True
    target_tf.margin_left = Inches(0.15)
    target_tf.margin_top = Inches(0.15)
    p_t = target_tf.paragraphs[0]
    p_t.text = "🎯 Core Verticals"
    p_t.font.bold = True
    p_t.font.size = Pt(13)
    p_t.font.color.rgb = WARM_BROWN
    p_t.font.name = 'Georgia'
    p_tb = target_tf.add_paragraph()
    p_tb.text = "Finance & Banking\nHealthcare EMR\nSupply Chain & POS\nEnterprise M&A Audit"
    p_tb.font.size = Pt(11)
    p_tb.font.color.rgb = DARK_INK
    p_tb.font.name = 'Georgia'
    p_tb.space_before = Pt(8)
    
    # 2 Large Detail Cards
    add_card(slide, Inches(3.3), Inches(1.8), Inches(4.5), Inches(4.9),
             "High-Impact Operational Savings",
             "• Regulatory Compliance: Auto-builds fully audit-compliant logs (Data Lineage) for regulatory bodies, avoiding multimillion-dollar reporting errors.\n\n"
             "• Labor Reductions: Reduces manual review timelines from weeks to seconds, allowing financial analysts to shift focus to core decision modeling and strategic insights.\n\n"
             "• Reduced Capital At Risk: Discovers silent internal data mismatches before they propagate into published quarterly financials or public M&A documents.",
             title_color=DEEP_BLUE, bg_color=CARD_BG)
             
    add_card(slide, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.9),
             "Commercial Monetization Pathways",
             "• Tiered SaaS License: Charged based on active user headcount, volume of active data pipelines, and database processing frequency.\n\n"
             "• Secure Enterprise API: Dedicated connectors designed to plug NeuralNexus directly into legacy corporate Business Intelligence tools like Tableau, Power BI, and ERPs.\n\n"
             "• Premium Compliance Exports: Custom billing model for exporting signed, cryptographically hashed audit trails for official external corporate auditing.",
             title_color=DEEP_BLUE, bg_color=CARD_BG)

    # ── SLIDE 11: COMPARATIVE USP ────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    add_title_text(slide, "Originality & Unique Selling Proposition")
    add_subtitle(slide, "Comparing NeuralNexus with conventional AI assistants, dashboard systems, and ETL tools.")
    add_bottom_bar(slide)
    
    headers_usp = ["Feature Capability", "Standard LLMs (ChatGPT)", "Traditional BI (Power BI)", "Legacy ETL (Informatica)", "NeuralNexus Fabric"]
    rows_usp = [
        ["Federated Ingestion", "❌ Manual upload only", "❌ Requires separate ETL", "✅ Yes, complex manual setup", "✅ Auto-federate multiple files"],
        ["Out-of-Box Anomaly Detection", "❌ None", "⚠️ Rule-based manual metrics", "⚠️ Requires complex configurations", "✅ Continuous autonomous scans"],
        ["Cross-Source Reconciliation", "❌ Fails, hallucinates facts", "❌ None", "❌ None", "✅ Trust-weighted reasoning"],
        ["Independent Answer Auditing", "❌ None", "❌ None", "❌ None", "✅ Five-dimension Critic QA gate"],
        ["Fact Provenance & Lineage", "❌ None", "⚠️ Restricted to DB source info", "✅ Data movement mapping only", "✅ Cell-level audit trail in answers"],
        ["Proactive Financial Alerts", "❌ Requires manual query prompt", "⚠️ Static reports/manual dashboards", "❌ None", "✅ Proactive background insight feeds"],
        ["Dynamic Knowledge Mapping", "❌ None", "❌ None", "❌ None", "✅ Auto-built NetworkX graph UI"]
    ]
    
    add_table_slide(slide, headers_usp, rows_usp, left=Inches(0.8), top=Inches(1.8), col_widths=[Inches(2.5), Inches(2.1), Inches(2.2), Inches(2.2), Inches(2.5)])

    # ── SLIDE 12: ROADMAP (VISUAL TIMELINE) ──────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    add_title_text(slide, "NeuralNexus Strategic Execution Roadmap")
    add_subtitle(slide, "A step-by-step scaling plan from current production MVP to full enterprise integrations.")
    add_bottom_bar(slide)
    
    # 5 step timeline
    circle_size = Inches(1.3)
    top_pos = Inches(2.5)
    
    steps = [
        ("Phase 1", "Current Production MVP", "5 agents integrated. React frontend + FastAPI backend. Streaming and KG active.", DEEP_BLUE),
        ("Phase 2", "Multi-Tenant Integration", "Implementing SSO/SAML, database auth, and granular user roles.", ACCENT_TEAL),
        ("Phase 3", "Live Connectors", "Connecting to PostgreSQL, Snowflake, and BigQuery for automated schedules.", WARM_BROWN),
        ("Phase 4", "Deep AI Governance", "Enterprise integration with Lobster Trap for complete prompt audit trails.", ACCENT_GOLD),
        ("Phase 5", "Agent Marketplace", "Allowing corporate teams to build custom task-specific analytical agents.", ACCENT_RED)
    ]
    
    for i, (phase, title, desc, color) in enumerate(steps):
        left = Inches(0.8) + i * Inches(2.4)
        
        # Circle
        add_icon_circle(slide, left + Inches(0.55), top_pos, circle_size, color, f"P{i+1}")
        
        # Text block
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top_pos + Inches(1.5), Inches(2.2), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = BORDER_COLOR
        box_tf = box.text_frame
        box_tf.word_wrap = True
        
        p1 = box_tf.paragraphs[0]
        p1.text = phase
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = color
        p1.font.name = 'Georgia'
        
        p2 = box_tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = DEEP_BLUE
        p2.font.name = 'Georgia'
        p2.space_before = Pt(4)
        
        p3 = box_tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(10)
        p3.font.color.rgb = DARK_INK
        p3.font.name = 'Georgia'
        p3.space_before = Pt(6)
        
        # Connecting arrow
        if i < 4:
            add_arrow_right(slide, left + Inches(2.1), top_pos + Inches(0.5), Inches(0.4), Inches(0.2))

    # ── SLIDE 13: TOOLS & TECHNOLOGIES ───────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    add_title_text(slide, "The Complete Technical Stack")
    add_subtitle(slide, "Standard, highly robust open-source tools and enterprise-ready framework structures.")
    add_bottom_bar(slide)
    
    # 3 Large categories side-by-side
    add_card(slide, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.9),
             "🚀 Robust Backend Ecosystem",
             "• Language: Python 3.13\n\n"
             "• Web Server: FastAPI async core, served via Uvicorn production container.\n\n"
             "• Core Agent Logic: Google Gemini 2.5 Flash Lite for agent reasoning.\n\n"
             "• Vector Store: FAISS (CPU) for nearest-neighbor semantic search indexing.\n\n"
             "• Data Processing: Pandas, NumPy, SciPy, and Scikit-learn.\n\n"
             "• Graph Processing: NetworkX backend, dynamically exported to PyVis HTML.",
             title_color=DEEP_BLUE, bg_color=CARD_BG)
             
    add_card(slide, Inches(4.7), Inches(1.8), Inches(3.6), Inches(4.9),
             "🎨 Visual Frontend Interface",
             "• Framework: React 18 component architecture, built with lightning-fast Vite.\n\n"
             "• UI Styling: Clean, premium Vanilla CSS library custom-coded for a beautiful 'Bookish' theme.\n\n"
             "• Icons System: Lucide React high-resolution graphics library.\n\n"
             "• Real-Time Streams: Customized Fetch Streams API to process and decode token NDJSON payloads dynamically.",
             title_color=DEEP_BLUE, bg_color=CARD_BG)
             
    add_card(slide, Inches(8.6), Inches(1.8), Inches(3.9), Inches(4.9),
             "🔧 Deployment & Lifecycle Tools",
             "• Frontend Hosting: Vercel cloud architecture (connected to GitHub branch updates).\n\n"
             "• Backend Deployment: Render cloud service with secure virtual environments.\n\n"
             "• Code Repository: Managed securely under GitHub version control.\n\n"
             "• Secret Protection: Multi-layered environment variables via python-dotenv (.env.example).",
             title_color=DEEP_BLUE, bg_color=CARD_BG)

    # ── SLIDE 14: LIVE DEMO FLOW ─────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    add_title_text(slide, "Live Application Demonstration Guide")
    add_subtitle(slide, "Step-by-step user interaction walkthrough on our production-ready platform.")
    add_bottom_bar(slide)
    
    # Numbered cards flow (5 steps)
    step_w = Inches(2.2)
    step_spacing = Inches(0.15)
    
    demo_steps = [
        ("Step 1", "Launch App", "Open neuralnexus-omega.vercel.app. Direct public link, no Vercel login required.", DEEP_BLUE),
        ("Step 2", "Feed Datasets", "Click 'Init Demo Data' in the left menu to dynamically load 5 pre-built enterprise files.", ACCENT_TEAL),
        ("Step 3", "Run Pipeline", "Click 'Analyze' in the sidebar to execute pipeline scoring, contradictions, and graph building.", WARM_BROWN),
        ("Step 4", "Interactive Chat", "Type questions like: 'Summarize NA revenue' in the Query panel and watch NDJSON stream real-time.", ACCENT_GOLD),
        ("Step 5", "Explore UI", "Explore the dashboard tabs: read Quality scores, see conflicts, and explore the interactive graph.", ACCENT_GREEN)
    ]
    
    for i, (step, title, text, color) in enumerate(demo_steps):
        left = Inches(0.8) + i * (step_w + step_spacing)
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), step_w, Inches(4.9))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        box.line.width = Pt(2)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.2)
        tf.margin_left = Inches(0.15)
        
        p1 = tf.paragraphs[0]
        p1.text = step
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = color
        p1.font.name = 'Georgia'
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = DEEP_BLUE
        p2.font.name = 'Georgia'
        p2.space_before = Pt(6)
        
        p3 = tf.add_paragraph()
        p3.text = text
        p3.font.size = Pt(11)
        p3.font.color.rgb = DARK_INK
        p3.font.name = 'Georgia'
        p3.space_before = Pt(12)

    # ── SLIDE 15: CONCLUSION ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, PARCHMENT)
    add_title_text(slide, "NeuralNexus: Intelligence with Integrity")
    add_subtitle(slide, "Redefining how enterprises validate, reconcile, and query multi-source financial data.")
    add_bottom_bar(slide)
    
    # 4 Large summary cards
    add_card(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3),
             "Coordinated 5-Agent Architecture",
             "A sophisticated multi-agent pipeline using Gemini 2.5 Flash Lite for fast, low-latency analytical reasoning. Features continuous background validation and strict answer auditing by the Critic Agent.",
             title_color=DEEP_BLUE)
             
    add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.3),
             "Complete Factual Integrity",
             "Auto-detects numerical and semantic contradictions between data sources. Traces cell-level and text-span-level data lineage, completely neutralizing standard LLM hallucinations.",
             title_color=DEEP_BLUE)
             
    add_card(slide, Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.3),
             "Widescreen Interface & Graph",
             "Beautiful custom-styled Bookish interface featuring real-time streaming, interactive quality score dashboards, proactive insights, and dynamic PyVis knowledge graphs.",
             title_color=DEEP_BLUE)
             
    add_card(slide, Inches(6.8), Inches(4.4), Inches(5.7), Inches(2.3),
             "Hackathon Ready Deployment",
             "100% production-ready and fully public deployment on Vercel and Render. Pushed directly to GitHub. Connect to your enterprise database to analyze data with total integrity.",
             title_color=DEEP_BLUE)
             
    # Bold final callout
    callout = slide.shapes.add_textbox(Inches(0.8), Inches(6.75), Inches(11.7), Inches(0.4))
    callout_p = callout.text_frame.paragraphs[0]
    callout_p.text = "“NeuralNexus doesn't just answer questions — it questions the answers.”"
    callout_p.font.size = Pt(14)
    callout_p.font.bold = True
    callout_p.font.italic = True
    callout_p.font.color.rgb = WARM_BROWN
    callout_p.font.name = 'Georgia'
    callout_p.alignment = PP_ALIGN.CENTER
    
    # Save the polished masterpiece
    prs.save("NeuralNexus_Visual_Presentation.pptx")
    print("Masterpiece visual presentation saved successfully!")

if __name__ == "__main__":
    build_presentation()
