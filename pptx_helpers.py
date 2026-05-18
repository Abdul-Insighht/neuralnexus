"""Helper functions for visually rich PPTX generation."""
import collections, collections.abc
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Bookish color palette
PARCHMENT = RGBColor(250, 245, 235)
DARK_INK = RGBColor(30, 40, 50)
WARM_BROWN = RGBColor(120, 80, 40)
DEEP_BLUE = RGBColor(30, 58, 95)
ACCENT_GOLD = RGBColor(180, 140, 60)
ACCENT_RED = RGBColor(180, 60, 60)
ACCENT_GREEN = RGBColor(40, 120, 80)
ACCENT_TEAL = RGBColor(40, 100, 120)
LIGHT_BOX = RGBColor(240, 235, 220)
WHITE = RGBColor(255, 255, 255)
CARD_BG = RGBColor(245, 240, 228)
BORDER_COLOR = RGBColor(200, 185, 160)

def set_bg(slide, color=PARCHMENT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_decorative_bar(slide, left=0, top=0, width=Inches(13.33), height=Inches(0.06), color=ACCENT_GOLD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_bottom_bar(slide):
    add_decorative_bar(slide, left=0, top=Inches(7.2), width=Inches(13.33), height=Inches(0.3), color=DEEP_BLUE)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(12), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "NeuralNexus — Track 4: Data & Intelligence — AI Agent Olympics 2026"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.name = 'Georgia'
    p.alignment = PP_ALIGN.CENTER

def add_card(slide, left, top, width, height, title, body, title_color=DEEP_BLUE, bg_color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = BORDER_COLOR
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = 'Georgia'
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK_INK
    p2.font.name = 'Georgia'
    p2.space_before = Pt(6)
    return shape

def add_title_text(slide, text, left=Inches(0.8), top=Inches(0.4), width=Inches(11.5), size=Pt(32), color=DEEP_BLUE):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Georgia'
    return tb

def add_subtitle(slide, text, left=Inches(0.8), top=Inches(1.1), width=Inches(11.5), size=Pt(16), color=WARM_BROWN):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.italic = True
    p.font.color.rgb = color
    p.font.name = 'Georgia'
    return tb

def add_body_text(slide, text, left=Inches(0.8), top=Inches(1.8), width=Inches(11.5), height=Inches(5), size=Pt(14)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = DARK_INK
        p.font.name = 'Georgia'
        p.space_after = Pt(4)
    return tb

def add_icon_circle(slide, left, top, size, color, label):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Georgia'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

def add_arrow_right(slide, left, top, width=Inches(0.5), height=Inches(0.3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GOLD
    shape.line.fill.background()
    return shape

def add_arrow_down(slide, left, top, width=Inches(0.3), height=Inches(0.4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GOLD
    shape.line.fill.background()
    return shape

def add_stat_box(slide, left, top, number, label, color=DEEP_BLUE):
    w, h = Inches(2.2), Inches(1.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Georgia'
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(220, 215, 200)
    p2.font.name = 'Georgia'
    p2.alignment = PP_ALIGN.CENTER
    return shape

def add_table_slide(slide, headers, rows, left=Inches(0.8), top=Inches(2.0), col_widths=None):
    cols = len(headers)
    r = len(rows) + 1
    if col_widths is None:
        tw = Inches(11.5)
        col_widths = [tw // cols] * cols
    table_shape = slide.shapes.add_table(r, cols, left, top, sum(col_widths), Inches(0.4 * r))
    table = table_shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DEEP_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Georgia'
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if ri % 2 == 0 else PARCHMENT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_INK
                p.font.name = 'Georgia'
    return table_shape
