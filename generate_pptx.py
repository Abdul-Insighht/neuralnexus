import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_bookish_theme(slide):
    # Set background color to parchment
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 245, 235)

def format_title(title_shape, text):
    title_shape.text = text
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.name = 'Georgia'
        paragraph.font.color.rgb = RGBColor(30, 40, 50)
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True

def format_body(body_shape, text, is_title=False):
    body_shape.text = text
    for paragraph in body_shape.text_frame.paragraphs:
        paragraph.font.name = 'Georgia'
        paragraph.font.color.rgb = RGBColor(40, 50, 60)
        paragraph.font.size = Pt(24) if is_title else Pt(18)

def main():
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    slides_data = [
        {
            "title": "NeuralNexus\nAutonomous Financial Intelligence Fabric",
            "content": (
                "5 AI Agents That Detect Contradictions, Reconcile Conflicts, "
                "Forecast Risks, and Audit Every Answer — With Full Data Lineage.\n\n"
                "Track: Track 4 — Data & Intelligence\n"
                "Hackathon: AI Agent Olympics 2026 — lablab.ai\n"
                "Team Members:\n"
                "  - Hafiz Abdul Rehman\n"
                "  - Meer Hamza\n"
                "  - Sayyam Akram\n\n"
                "Live App: neuralnexus-omega.vercel.app\n"
                "GitHub: github.com/Abdul-Insighht/neuralnexus"
            )
        },
        {
            "title": "The Problem: $12.9M Lost Annually to Poor Data",
            "content": (
                "• Data Silos: Financial data lives across 5–15 disconnected systems.\n"
                "• Silent Contradictions: E.g., Sales DB vs Financial Report. No one catches it until audit.\n"
                "• Reactive Analysis: Teams spend 80% finding problems, not solving them.\n"
                "• Zero Traceability: Hard to trace numbers back to sources.\n"
                "• Unaudited AI Answers: Chat tools give answers without citing sources or uncertainty.\n\n"
                "Result: Erroneous decisions, compliance violations, regulatory fines."
            )
        },
        {
            "title": "Track 4 Alignment — Data & Intelligence",
            "content": (
                "How NeuralNexus Addresses Track 4:\n"
                "• RAG over multi-source data: Temporal Vector Store with Gemini Embeddings over federated data.\n"
                "• AI-powered data validation: ValidatorAgent runs autonomous anomaly scans (Z-score, IQR).\n"
                "• Analytics agents for natural language querying: OrchestratorAgent with real-time streaming.\n"
                "• Anomaly detection and forecasting: ValidatorAgent (stats) + ForecasterAgent (trends, risks).\n"
                "• Knowledge graph extraction: DynamicKnowledgeGraph auto-extracts entities with contradiction edges."
            )
        },
        {
            "title": "Our Solution: An Autonomous Intelligence Fabric",
            "content": (
                "NeuralNexus is an autonomous multi-agent intelligence fabric that:\n"
                "1. Federates data from any format (CSV, Excel, JSON, PDF, Images).\n"
                "2. Automatically detects contradictions across all sources.\n"
                "3. Reconciles conflicts using AI reasoning with trust-weighted analysis.\n"
                "4. Forecasts trends and risks before they become problems.\n"
                "5. Audits every AI answer for accuracy, completeness, and bias.\n"
                "6. Traces every fact back to its exact source, row, and column.\n\n"
                "It watches your data 24/7, catches what humans miss, and proves every number it cites."
            )
        },
        {
            "title": "System Architecture",
            "content": (
                "Multi-Agent Orchestration Pipeline:\n"
                "• OrchestratorAgent (The Brain): Routes, coordinates, synthesizes.\n"
                "• ValidatorAgent (The Watchdog): Scans for anomalies and drift.\n"
                "• ReconcilerAgent (The Judge): Resolves cross-source conflicts.\n"
                "• ForecasterAgent (The Oracle): Predicts trends and risks.\n"
                "• CriticAgent (The Auditor): QA review before user sees the answer.\n"
                "• Core Infrastructure: Knowledge Graph, Vector Store, Lineage Tracker."
            )
        },
        {
            "title": "The 5 AI Agents",
            "content": (
                "1. OrchestratorAgent: Central coordinator, routes queries, streams NDJSON real-time.\n"
                "2. ValidatorAgent: Anomaly detection (Z-Score, IQR, temporal drift, sudden spikes).\n"
                "3. ReconcilerAgent: Resolves contradictions using trust-weighted analysis.\n"
                "4. ForecasterAgent: Predictive intelligence, trend direction, seasonality, risk scoring.\n"
                "5. CriticAgent: Quality assurance gate. Scores Accuracy, Completeness, Actionability. Adds caveats."
            )
        },
        {
            "title": "Core Infrastructure Modules",
            "content": (
                "• Data Federation Hub: Multi-format ingestion, auto-schema inference, data fingerprinting.\n"
                "• Data Quality Scorer: Completeness, Consistency, Timeliness. Grades A-F.\n"
                "• Contradiction Detector: Structured vs Structured, Structured vs Text, Text vs Text.\n"
                "• Data Lineage Tracker: Full provenance for every fact (row, column, transform step).\n"
                "• Dynamic Knowledge Graph: Entities and temporal/contradiction edges via PyVis.\n"
                "• Temporal Vector Store: Semantic search with recency-weighted retrieval."
            )
        },
        {
            "title": "User Interface — Dashboard Modules",
            "content": (
                "• Intelligence Query: Chat interface with real-time streaming, citations, and confidence scores.\n"
                "• Quality Dashboard: Visual breakdown of quality grades (A-F) and actionable issues.\n"
                "• Contradictions: Lists detected conflicts with severity, entity, and exact values.\n"
                "• Proactive Insights: AI-generated risk briefings without user prompting.\n"
                "• Knowledge Graph: Interactive visual graph showing entity relationships.\n"
                "• Data Lineage: Full audit trail for every fact."
            )
        },
        {
            "title": "Gemini Integration",
            "content": (
                "• Gemini 2.5 Flash Lite: Fast, low-latency reasoning for all 5 agents.\n"
                "• Gemini Embedding-2 (768-dim): For semantic search in Vector Store.\n"
                "• Gemini Multimodal Vision: For PDF and Image content extraction.\n"
                "• Real-Time Streaming: Uses Gemini stream=True API for token-by-token UI feedback.\n"
                "• Rate Limiting: Thread-safe global limiter prevents quota errors."
            )
        },
        {
            "title": "Business Value & Market Scope",
            "content": (
                "• Target Industries: Finance & Banking, Healthcare, Supply Chain, Gov/Compliance, M&A.\n"
                "• Total Addressable Market (TAM): $15.8B (Enterprise Data Quality Management).\n"
                "• Serviceable Addressable Market (SAM): $4.2B (AI Data Reconciliation).\n"
                "• Revenue Streams:\n"
                "  1. SaaS Subscription (per-user/per-source)\n"
                "  2. Enterprise API Licensing\n"
                "  3. Compliance audit trail exports."
            )
        },
        {
            "title": "Originality & USP",
            "content": (
                "What Makes NeuralNexus Different?\n"
                "• ChatGPT: General Q&A, but no multi-source data, contradiction detection, or lineage.\n"
                "• Tableau / BI Tools: Manual data prep, no AI reasoning or auto-reconciliation.\n"
                "• Informatica / Talend: Expensive, no natural language interface.\n"
                "• NeuralNexus USP: The only tool combining multi-source federation, autonomous contradiction detection, "
                "AI reconciliation, answer auditing, and full data lineage in a single fabric."
            )
        },
        {
            "title": "Future Roadmap & Scalability",
            "content": (
                "• Phase 1 (Current): MVP with 5 agents, deployed on Vercel + Render.\n"
                "• Phase 2: Multi-tenant support, user auth, RBAC.\n"
                "• Phase 3: Real-time database connectors (PostgreSQL, Snowflake).\n"
                "• Phase 4: Enterprise Lobster Trap integration for AI governance.\n"
                "• Phase 5: Custom agent marketplace.\n"
                "• Scalability: Stateless agents scale horizontally. Swap Flash Lite for Pro for higher reasoning tasks."
            )
        },
        {
            "title": "Tools & Technologies",
            "content": (
                "Backend: Python 3.13, FastAPI, Uvicorn, Google Gemini API, Pandas, NumPy, NetworkX, PyVis, FAISS, PyMuPDF, Pydantic.\n"
                "Frontend: React 18, Vite, Vanilla CSS (Bookish theme), Fetch Streams API.\n"
                "DevOps: Vercel (Frontend), Render (Backend), GitHub, dotenv."
            )
        },
        {
            "title": "Live Demo Flow",
            "content": (
                "1. Open App: neuralnexus-omega.vercel.app\n"
                "2. Init Demo Data: Loads 5 enterprise sources.\n"
                "3. Analyze: Runs full pipeline (quality, anomalies, KG).\n"
                "4. Query: Ask a question, watch real-time stream with citations.\n"
                "5. Dashboards: View Quality, Contradictions, Insights, Knowledge Graph, and Data Lineage."
            )
        },
        {
            "title": "Conclusion",
            "content": (
                "NeuralNexus Solves a $12.9B Problem:\n"
                "• 5 Autonomous AI Agents working together.\n"
                "• Real-time streaming for instant feedback.\n"
                "• Automatic contradiction detection & AI reconciliation.\n"
                "• Proactive insights and full data lineage.\n"
                "• Production ready on Vercel + Render.\n\n"
                "NeuralNexus doesn't just answer questions — it questions the answers."
            )
        }
    ]
    
    for i, slide_data in enumerate(slides_data):
        layout = title_slide_layout if i == 0 else bullet_slide_layout
        slide = prs.slides.add_slide(layout)
        apply_bookish_theme(slide)
        
        format_title(slide.shapes.title, slide_data['title'])
        if len(slide.placeholders) > 1:
            format_body(slide.placeholders[1], slide_data['content'], is_title=(i==0))
            
    prs.save('NeuralNexus_Presentation.pptx')
    print("Presentation created successfully as NeuralNexus_Presentation.pptx")

if __name__ == '__main__':
    main()
